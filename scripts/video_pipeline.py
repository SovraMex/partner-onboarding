"""Minimal video pipeline: pptx + chunks -> .mp4 with Dalia voice.
Per-slide duration == its chunk audio duration. No Ken Burns, no watermark."""
import asyncio
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

import edge_tts
import imageio_ffmpeg
import pythoncom
import win32com.client
from pptx import Presentation

FFMPEG = imageio_ffmpeg.get_ffmpeg_exe()


def extract_slide_texts(pptx_path):
    """Return list of slide text content (joined from all shapes)."""
    prs = Presentation(str(pptx_path))
    out = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            parts.append(run.text)
        out.append(' '.join(parts).strip())
    return out


def pptx_to_pngs(pptx_path, out_dir):
    """Render each slide to a PNG via PowerPoint COM. Returns sorted list of PNG paths."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    pythoncom.CoInitialize()
    powerpoint = None
    pres = None
    try:
        powerpoint = win32com.client.Dispatch('PowerPoint.Application')
        # Open hidden
        pres = powerpoint.Presentations.Open(
            str(Path(pptx_path).absolute()),
            ReadOnly=True,
            Untitled=False,
            WithWindow=False
        )
        # Export to PNG. PowerPoint's Export with FilterName='PNG' writes 1920x1080 by default
        # for 16:9 decks. Each slide goes as Slide1.PNG, Slide2.PNG, ...
        pres.SaveAs(str(out_dir.absolute()), 18)  # 18 = ppSaveAsPNG, writes one PNG per slide
        pres.Close()
        powerpoint.Quit()
    except Exception:
        if pres is not None:
            try: pres.Close()
            except Exception: pass
        if powerpoint is not None:
            try: powerpoint.Quit()
            except Exception: pass
        raise
    finally:
        pythoncom.CoUninitialize()

    # ppSaveAsPNG writes Slide1.PNG, Slide2.PNG, ... but on some versions it writes to a subfolder
    # with the same name as the deck. Handle both.
    pngs = sorted(out_dir.rglob('Slide*.PNG')) or sorted(out_dir.rglob('Slide*.png')) \
        or sorted(out_dir.rglob('*.PNG')) or sorted(out_dir.rglob('*.png'))
    # Filter to keep only the Slide-numbered ones and sort numerically
    import re as _re
    def slide_num(p):
        m = _re.search(r'Slide(\d+)', p.name)
        return int(m.group(1)) if m else 9999
    pngs = sorted([p for p in pngs if _re.search(r'Slide\d+', p.name)], key=slide_num)
    return pngs


async def generate_audio_async(text, voice, out_path):
    """Generate TTS via Edge. Returns duration in seconds."""
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(str(out_path))
    # Measure duration with ffprobe (via ffmpeg)
    proc = subprocess.run(
        [FFMPEG, '-i', str(out_path)],
        capture_output=True, text=True, errors='replace'
    )
    # ffmpeg writes duration to stderr like "Duration: HH:MM:SS.MS"
    import re as _re
    m = _re.search(r'Duration:\s*(\d+):(\d+):(\d+\.\d+)', proc.stderr)
    if not m:
        raise RuntimeError(f'Could not parse duration for {out_path}')
    h, mn, s = int(m.group(1)), int(m.group(2)), float(m.group(3))
    return h * 3600 + mn * 60 + s


def generate_audio(text, voice, out_path):
    return asyncio.run(generate_audio_async(text, voice, out_path))


def concat_audio(audio_paths, out_path):
    """Concatenate mp3 files into one via ffmpeg concat demuxer."""
    list_file = Path(out_path).parent / 'audio-list.txt'
    list_file.write_text('\n'.join(f"file '{Path(p).absolute().as_posix()}'" for p in audio_paths), encoding='utf-8')
    subprocess.run([
        FFMPEG, '-y', '-hide_banner', '-loglevel', 'error',
        '-f', 'concat', '-safe', '0', '-i', str(list_file),
        '-c', 'copy', str(out_path)
    ], check=True)


def assemble_video(png_paths, full_audio, slide_durations, output_path):
    """For each PNG, show it for slide_durations[i] seconds. Mux with full_audio.
    Output: H.264 mp4, 1920x1080, 30fps, slide duration EXACTLY matches each chunk's audio."""
    if len(png_paths) != len(slide_durations):
        raise RuntimeError(f'PNG count {len(png_paths)} != duration count {len(slide_durations)}')

    # Build input args: one image input per slide, with -loop 1 and -t duration
    args = [FFMPEG, '-y', '-hide_banner', '-loglevel', 'error']
    for png, dur in zip(png_paths, slide_durations):
        args += ['-loop', '1', '-t', f'{dur:.3f}', '-i', str(png)]
    # Audio input last
    args += ['-i', str(full_audio)]

    n = len(png_paths)
    # Build filter: scale + pad each slide to 1920x1080, then concat
    filter_parts = []
    for i in range(n):
        filter_parts.append(
            f'[{i}:v]scale=1920:1080:force_original_aspect_ratio=decrease,'
            f'pad=1920:1080:(ow-iw)/2:(oh-ih)/2:color=white,'
            f'setsar=1,fps=30,format=yuv420p[v{i}]'
        )
    concat_inputs = ''.join(f'[v{i}]' for i in range(n))
    filter_parts.append(f'{concat_inputs}concat=n={n}:v=1:a=0[outv]')
    filter_str = ';'.join(filter_parts)

    args += [
        '-filter_complex', filter_str,
        '-map', '[outv]',
        '-map', f'{n}:a',
        '-c:v', 'libx264',
        '-preset', 'medium',
        '-crf', '23',
        '-c:a', 'aac',
        '-b:a', '128k',
        '-shortest',
        str(output_path)
    ]
    subprocess.run(args, check=True)


# Cache stub used by external regen scripts via monkey-patch.
def _try_cache(slide_texts, pptx_stem):
    return None


def generate_video(pptx_path, script_text, output_path, voice='es-MX-DaliaNeural'):
    """End-to-end. Caller is expected to monkey-patch _try_cache to supply pre-computed chunks."""
    pptx_path = Path(pptx_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    work = Path(tempfile.mkdtemp(prefix='vidgen-'))
    try:
        # 1. Render pptx -> PNGs
        png_dir = work / 'pngs'
        pngs = pptx_to_pngs(pptx_path, png_dir)
        if not pngs:
            raise RuntimeError(f'No PNGs from {pptx_path}')
        num_slides = len(pngs)

        # 2. Get chunks from cache (caller monkey-patches _try_cache)
        slide_texts = extract_slide_texts(pptx_path)
        chunks = _try_cache(slide_texts, pptx_path.stem)
        if not chunks:
            raise RuntimeError(f'No cached matching for {pptx_path.stem} — cache miss')
        if len(chunks) != num_slides:
            raise RuntimeError(f'Cached chunks {len(chunks)} != slides {num_slides}')

        # 3. Generate audio per chunk + measure durations
        audio_paths = []
        durations = []
        for i, chunk in enumerate(chunks):
            audio_path = work / f'chunk-{i:02d}.mp3'
            dur = generate_audio(chunk, voice, audio_path)
            audio_paths.append(audio_path)
            durations.append(dur)

        # 4. Concat audio
        full_audio = work / 'narration.mp3'
        concat_audio(audio_paths, full_audio)

        # 5. Assemble video — slide i shown for durations[i] seconds
        assemble_video(pngs, full_audio, durations, output_path)
        return str(output_path)
    finally:
        shutil.rmtree(work, ignore_errors=True)
